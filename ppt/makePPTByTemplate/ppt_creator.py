# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import os
import warnings
from copy import deepcopy
import glob
import tempfile
import re


class PPTCreator:
    def __init__(self, template_dir="粉白简约"):
        """
        初始化PPT创建器

        参数:
            template_dir (str): 模板目录路径，默认为"粉白简约"
        """
        # 加载模板文件路径
        self.template_dir = template_dir
        self.title_template = os.path.join(template_dir, "title.pptx")
        self.contents_template = os.path.join(template_dir, "contents.pptx")  # 新增目录模板
        self.intro_template = os.path.join(template_dir, "intro.pptx")
        self.topic_template = os.path.join(template_dir, "topic.pptx")
        self.conclusion_template = os.path.join(template_dir, "conclusion.pptx")

        # 获取section目录下所有模板文件并按数字排序
        self.section_templates = sorted(
            glob.glob(os.path.join(template_dir, "section", "*.pptx")),
            key=lambda x: int(os.path.basename(x).split('.')[0]))

        if not self.section_templates:
            raise ValueError("在section目录中未找到任何模板文件")

        # 初始化模板计数器
        self.section_template_index = 0

        # 创建一个新的空白演示文稿
        self.prs = Presentation()

        # 从标题模板获取幻灯片大小并应用到新演示文稿
        title_prs = Presentation(self.title_template)
        self.prs.slide_width = title_prs.slide_width
        self.prs.slide_height = title_prs.slide_height

    def add_contents_slide(self, topics):
        """
        添加目录幻灯片（包含引言、各主题和总结）

        参数:
            topics (list): 主题名称列表
        """
        slide = self._copy_slide_from_template(self.contents_template)

        # 查找目录内容占位符
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "目录" in text:
                    # 保留"目录"标题
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.size = Pt(32)
                elif "内容" in text:
                    # 清空原有内容
                    shape.text_frame.clear()

                    # 添加固定项和带序号的目录项
                    contents = [
                        "1. 引言",  # 固定第一项
                        *[f"{i + 2}. {topic}" for i, topic in enumerate(topics)],  # 主题列表
                        f"{len(topics) + 2}. 总结"  # 固定最后一项
                    ]

                    for content in contents:
                        p = shape.text_frame.add_paragraph()
                        p.text = content
                        p.level = 0
                        p.font.size = Pt(24)
                        p.space_after = Pt(12)

                    # 设置自动调整文本框大小
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def _get_next_section_template(self):
        """
        获取下一个section模板文件路径，循环使用

        返回:
            str: 下一个模板文件路径
        """
        template_path = self.section_templates[self.section_template_index]
        self.section_template_index = (self.section_template_index + 1) % len(self.section_templates)
        return template_path

    def _copy_slide_from_template(self, template_path):
        """
        从模板中复制一张幻灯片到当前演示文稿，包括图片

        参数:
            template_path (str): 模板文件路径

        返回:
            Slide: 新创建的幻灯片对象
        """
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            template_prs = Presentation(template_path)
            template_slide = template_prs.slides[0]

            # 创建一个新的幻灯片布局
            blank_slide_layout = self.prs.slide_layouts[6]  # 6是空白布局
            new_slide = self.prs.slides.add_slide(blank_slide_layout)

            # 复制所有形状，包括图片
            for shape in template_slide.shapes:
                if hasattr(shape, 'image') and shape.image:  # 更安全的图片检测方式
                    try:
                        # 获取图片数据
                        image_bytes = shape.image.blob

                        # 创建临时文件来保存图片
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.jpeg') as temp_file:
                            temp_file.write(image_bytes)
                            temp_file_path = temp_file.name

                        # 添加图片到新幻灯片
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        new_slide.shapes.add_picture(
                            temp_file_path,
                            left, top, width, height
                        )

                        # 删除临时文件
                        os.unlink(temp_file_path)
                    except Exception as e:
                        print(f"无法复制图片: {e}")
                        continue
                else:
                    # 复制其他形状
                    try:
                        new_shape = deepcopy(shape)
                        new_slide.shapes._spTree.insert_element_before(new_shape.element, 'p:extLst')
                    except Exception as e:
                        print(f"无法复制形状: {e}")
                        continue

            # # 处理背景图片
            # try:
            #     if hasattr(template_slide.background, 'fill') and template_slide.background.fill.type == 1:
            #         background_bytes = template_slide.background.fill._pic.blob
            #         with tempfile.NamedTemporaryFile(delete=False, suffix='.jpeg') as temp_file:
            #             temp_file.write(background_bytes)
            #             temp_file_path = temp_file.name
            #
            #         new_slide.background.fill.solid()
            #         new_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            #         new_slide.background.fill.picture()
            #         new_slide.background.fill._pic.blob = background_bytes
            #
            #         os.unlink(temp_file_path)
            # except Exception as e:
            #     print(f"无法复制背景图片: {e}")

        self._remove_default_placeholders(new_slide)
        return new_slide

    def _remove_default_placeholders(self, slide):
        """
        移除幻灯片中的默认占位符文本

        参数:
            slide (Slide): 幻灯片对象
        """
        default_texts = [
            "单击此处添加标题", "Click to add title",
            "单击此处添加文本", "Click to add text",
            "单击此处添加副标题", "Click to add subtitle",
            "双击此处添加标题", "Double-click to add title"
        ]

        # 先收集要删除的形状
        shapes_to_remove = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                # 检查是否是默认占位符
                if any(default_text in shape.text_frame.text for default_text in default_texts):
                    shapes_to_remove.append(shape)
                # 清空只包含默认文本的段落
                for paragraph in shape.text_frame.paragraphs:
                    if any(default_text in paragraph.text for default_text in default_texts):
                        paragraph.text = ""

        # 删除形状
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def add_title_slide(self, title, subtitle=None):
        """
        添加标题幻灯片

        参数:
            title (str): 主标题文本
            subtitle (str, optional): 副标题文本，默认为None
        """
        slide = self._copy_slide_from_template(self.title_template)

        # 查找并填充标题
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "标题" in text:
                    # 设置标题文本
                    shape.text_frame.text = title
                    # 设置字体样式
                    shape.text_frame.paragraphs[0].font.size = Pt(36)
                    shape.text_frame.paragraphs[0].font.bold = True
                    # 确保文本居中
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                    # 移除副标题形状
                    if "副标题" in text and subtitle is None:
                        sp = shape._element
                        sp.getparent().remove(sp)
                    elif "副标题" in text and subtitle is not None:
                        shape.text_frame.text = subtitle
                        shape.text_frame.paragraphs[0].font.size = Pt(24)

    def add_intro_slide(self, content):
        """
        添加介绍幻灯片

        参数:
            content (str): 介绍内容文本
        """
        slide = self._copy_slide_from_template(self.intro_template)

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "引文" in text:
                    # 保留"引文"标题
                    shape.text_frame.paragraphs[0].font.bold = True
                elif "内容" in text:
                    shape.text_frame.text = content
                    shape.text_frame.paragraphs[0].font.size = Pt(18)
                    # 设置自动调整文本框大小
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def add_topic_slide(self, topic_name):
        """
        添加主题幻灯片

        参数:
            topic_name (str): 主题名称
        """
        slide = self._copy_slide_from_template(self.topic_template)

        # 查找并填充主题名称
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "主题" in text:
                    shape.text_frame.text = topic_name
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.size = Pt(32)
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER

    def add_subtopic_slide(self, subtopic):
        """
        添加子主题幻灯片 (相当于原来的section)

        参数:
            subtopic (dict): 包含子主题信息的字典
        """
        template_path = self._get_next_section_template()
        slide = self._copy_slide_from_template(template_path)

        # 将subtopic转换为与原来section相同的格式
        section_data = {
            "section_title": subtopic["name"],
            "subsections": [
                {
                    "subtitle": section["title"],
                    "content": section["content"]
                } for section in subtopic["sections"]["each_section"]
            ]
        }

        self._fill_simple_layout(slide, section_data)

    def _fill_simple_layout(self, slide, section):
        """
        填充简单布局内容

        参数:
            slide (Slide): 幻灯片对象
            section (dict): 包含章节标题和子章节的字典
        """
        # 首先查找并替换章节标题
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "章节标题" in text:
                    shape.text_frame.text = section["section_title"]
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.size = Pt(28)

        # 处理各个小节
        subsections = section["subsections"]
        for i in range(1, 4):  # 支持最多3个小节
            if i > len(subsections):
                break

            subsection = subsections[i - 1]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if f"小标题{i}" in text or f"子标题{i}" in text:  # 兼容新旧模板
                        shape.text_frame.text = subsection["subtitle"]
                        shape.text_frame.paragraphs[0].font.bold = True
                        shape.text_frame.paragraphs[0].font.size = Pt(20)
                    elif f"内容{i}" in text:
                        shape.text_frame.text = subsection["content"]
                        shape.text_frame.paragraphs[0].font.size = Pt(16)
                        # 设置自动调整文本框大小
                        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def add_conclusion_slide(self, conclusion_sections):
        """
        添加总结幻灯片

        参数:
            conclusion_sections (list): 包含总结部分的字典列表
        """
        slide = self._copy_slide_from_template(self.conclusion_template)

        # 将总结部分转换为与原来section相同的格式
        section_data = {
            "section_title": "总结",
            "subsections": [
                {
                    "subtitle": section["title"],
                    "content": section["content"]
                } for section in conclusion_sections
            ]
        }

        self._fill_simple_layout(slide, section_data)

    def save(self, filename):
        """
        保存PPT文件

        参数:
            filename (str): 要保存的文件名
        """
        # 忽略保存时的警告
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            self.prs.save(filename)



# # 使用示例
# data =
#
#
# creator = PPTCreator(template_dir="粉白简约")
#
# # 1. 添加标题页 - 从字典中提取description
# creator.add_title_slide(data["ppt_title"]["description"])
#
# # 2. 添加目录页
# topic_names = [topic["name"] for topic in data["structure"]["topics"]["each_topic"]]
# creator.add_contents_slide(topic_names)
#
# # 3. 添加引文页
# creator.add_intro_slide(data["structure"]["introduction"]["description"])
#
# # 4. 添加主题页和子主题页
# for topic in data["structure"]["topics"]["each_topic"]:
#     # 添加主题页
#     creator.add_topic_slide(topic["name"])
#
#     # 添加子主题页
#     for subtopic in topic["subtopics"]["each_subtopic"]:
#         creator.add_subtopic_slide(subtopic)
#
# # 5. 添加总结页
# creator.add_conclusion_slide(data["structure"]["conclusion"]["sections"]["each_section"])
#
# # 保存PPT - 使用ppt_title作为文件名
# ppt_title = data["ppt_title"]["description"]
# # 移除Windows文件名中不允许的特殊字符
# safe_filename = re.sub(r'[\\/*?:"<>|]', "", ppt_title)
# output_filename = f"{safe_filename}.pptx"
#
# creator.save(output_filename)
# print(f"PPT生成完成，已保存为: {output_filename}")